import json
import re
from io import BytesIO
from typing import Any, Dict, List, Optional
from uuid import UUID

from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from sqlalchemy.orm import Session

from langchain_openai import OpenAIEmbeddings

from app.core.config import settings
from app.models import Agent, Embedding
from app.core.logging import logger


class EmbeddingService:
    def __init__(self, db: Session):
        self.db = db
        self.embedding_client = OpenAIEmbeddings(
            api_key=settings.OPENAI_API_KEY,
            model="text-embedding-3-small",
        )

    async def ingest_file(
        self,
        agent: Agent,
        filename: str,
        content_type: Optional[str],
        data: bytes,
        *,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
        batch_size: Optional[int] = None,
    ) -> Dict[str, Any]:
        extension = self._determine_extension(filename, content_type)
        text = self._extract_text(extension, data, filename)
        cleaned = self._clean_text(text)
        chunk_size_val = chunk_size if chunk_size and chunk_size > 0 else 500
        overlap_val = chunk_overlap if chunk_overlap and chunk_overlap >= 0 else 100
        batch_size_val = batch_size if batch_size and batch_size > 0 else 64

        char_count = len(cleaned)
        adjusted = False
        if char_count > 500_000:
            chunk_size_val = min(chunk_size_val, 200)
            overlap_val = min(overlap_val, 40)
            batch_size_val = min(batch_size_val, 24)
            adjusted = True
        elif char_count > 250_000:
            chunk_size_val = min(chunk_size_val, 300)
            overlap_val = min(overlap_val, 60)
            batch_size_val = min(batch_size_val, 32)
            adjusted = True

        if overlap_val >= chunk_size_val:
            overlap_val = max(0, chunk_size_val // 4)

        chunks = self._chunk_text(cleaned, chunk_size=chunk_size_val, overlap=overlap_val)

        if not chunks:
            raise ValueError("File did not contain readable text after cleaning.")

        if adjusted:
            logger.info(
                "Large document detected; adjusted chunk parameters",
                char_count=char_count,
                chunk_size=chunk_size_val,
                overlap=overlap_val,
                batch_size=batch_size_val,
            )

        vectors = self._embed_in_batches(chunks, batch_size_val)

        records: List[Embedding] = []
        total_chunks = len(chunks)
        for index, (chunk, vector) in enumerate(zip(chunks, vectors)):
            metadata = {
                "source": filename,
                "chunk": index,
                "total_chunks": total_chunks,
                "content_type": content_type,
            }
            record = Embedding(
                agent_id=agent.id,
                content=chunk,
                embedding=vector,
                metadata_json=metadata,
            )
            self.db.add(record)
            records.append(record)

        self.db.flush()
        ids = [str(record.id) for record in records]
        self.db.commit()

        return {
            "chunks": total_chunks,
            "embedding_ids": ids,
            "chunk_size": chunk_size_val,
            "chunk_overlap": overlap_val,
            "batch_size": batch_size_val,
            "characters": char_count,
        }

    def _determine_extension(self, filename: str, content_type: Optional[str]) -> str:
        if filename.lower().endswith(".pdf") or content_type == "application/pdf":
            return "pdf"
        if filename.lower().endswith(".docx") or content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return "docx"
        if filename.lower().endswith(".pptx") or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return "pptx"
        if filename.lower().endswith(".txt") or content_type == "text/plain":
            return "txt"
        raise ValueError("Unsupported file type. Provide pdf, docx, pptx, or txt.")

    def _extract_text(self, extension: str, data: bytes, filename: str) -> str:
        if extension == "txt":
            return data.decode("utf-8", errors="ignore")

        buffer = BytesIO(data)

        if extension == "pdf":
            reader = PdfReader(buffer)
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n".join(pages)

        if extension == "docx":
            document = Document(buffer)
            paragraphs = [para.text for para in document.paragraphs]
            return "\n".join(paragraphs)

        if extension == "pptx":
            presentation = Presentation(buffer)
            texts: List[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)

        raise ValueError(f"Unsupported file extension for {filename}")

    def _clean_text(self, text: str) -> str:
        text = text.replace("\x00", " ")
        text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", " ", text)
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    def _chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
        words = text.split()
        if not words:
            return []

        if len(words) <= chunk_size:
            return [" ".join(words)]

        chunks: List[str] = []
        start = 0
        while start < len(words):
            end = min(len(words), start + chunk_size)
            chunk = " ".join(words[start:end])
            chunks.append(chunk)
            if end == len(words):
                break
            start = max(end - overlap, start + 1)
        return chunks

    def _embed_in_batches(self, chunks: List[str], batch_size: int) -> List[List[float]]:
        vectors: List[List[float]] = []
        max_tokens_per_batch = 250_000
        start = 0
        while start < len(chunks):
            end = min(len(chunks), start + batch_size)
            batch = chunks[start:end]

            # adjust batch size if necessary to respect token limits
            while batch:
                estimated_tokens = sum(max(1, len(item) // 4) for item in batch)
                if estimated_tokens <= max_tokens_per_batch:
                    break
                if len(batch) == 1:
                    break
                end = start + max(1, len(batch) // 2)
                batch = chunks[start:end]

            vectors.extend(self.embedding_client.embed_documents(batch))
            start += len(batch)
        return vectors

    def get_relevant_chunks(
        self,
        agent_id: UUID,
        query: str,
        top_k: int = 3,
    ) -> List[Dict[str, Any]]:
        if not query.strip():
            return []

        query_vector = self.embedding_client.embed_query(query)

        distance = Embedding.embedding.cosine_distance(query_vector)

        rows = (
            self.db.query(
                Embedding.content,
                Embedding.metadata_json,
                distance.label("distance"),
            )
            .filter(Embedding.agent_id == agent_id)
            .order_by(distance)
            .limit(top_k)
            .all()
        )

        return [
            {
                "content": row[0],
                "metadata": row[1] or {},
                "distance": float(row[2]) if row[2] is not None else None,
            }
            for row in rows
        ]
